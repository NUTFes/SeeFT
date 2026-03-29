#!/usr/bin/env python3
"""Convert 配線マニュアル docx content into a well-structured portrait pptx for mobile."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)
IMG = os.path.join(SCRIPT_DIR, "docx_images")
OUT = os.path.join(PROJECT_ROOT, "docs", "manuals", "01_44th_配線マニュアル.pptx")

prs = Presentation()
# Portrait 9:16 for mobile
prs.slide_width = Inches(7.5)
prs.slide_height = Inches(13.333)

W = 7.5  # slide width in inches
MARGIN = 0.4
CONTENT_W = W - MARGIN * 2  # 6.7

# Color scheme
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE_BG = RGBColor(0xD6, 0xE4, 0xF0)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
TABLE_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)
WARN_BG = RGBColor(0xFF, 0xF3, 0xCD)

def add_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, text, y=Inches(0.2), height=Inches(0.7)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y, prs.slide_width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_text_box(slide, text, left, top, width, height, font_size=16, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, items, left, top, width, height, font_size=16):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "● " + item if not item.startswith("　") else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
    return txBox

def add_image_safe(slide, img_name, left, top, width=None, height=None):
    path = os.path.join(IMG, img_name)
    if os.path.exists(path):
        kwargs = {"image_file": path, "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        return slide.shapes.add_picture(**kwargs)
    return None

def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.4)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows)
    table = table_shape.table
    col_w = int(width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_w
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = BLACK
                p.alignment = PP_ALIGN.CENTER
    return table_shape

def add_accent_box(slide, text, left, top, width, height, bg_color=LIGHT_BLUE_BG, font_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = BLACK
    return shape

L = Inches(MARGIN)  # left margin
CW = Inches(CONTENT_W)  # content width

# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text_box(slide, "配線マニュアル", L, Inches(4), CW, Inches(1.5),
             font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, "第44回 技大祭", L, Inches(5.5), CW, Inches(0.8),
             font_size=24, color=RGBColor(0xA0, 0xC4, 0xE8), alignment=PP_ALIGN.CENTER)
add_text_box(slide, "準備日 【電力配線】シフト該当者向け", L, Inches(6.3), CW, Inches(0.8),
             font_size=18, color=RGBColor(0xA0, 0xC4, 0xE8), alignment=PP_ALIGN.CENTER)

# ===== SLIDE 2: 該当シフト & 集合場所 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "該当シフト・集合情報")

add_accent_box(slide, "準備日の【電力配線】のシフト該当者は、\nこの資料を参考に業務にあたること。",
               L, Inches(1.2), CW, Inches(0.8), font_size=15)

info_items = [
    ("集合場所", "B講義室"),
    ("集合時間", "17:50"),
    ("開始時間", "18:00"),
    ("終了時間", "22:00"),
    ("所要時間", "配線 : 4時間"),
]
for i, (label, value) in enumerate(info_items):
    col = i % 2
    row = i // 2
    x = Inches(MARGIN + col * 3.5)
    y = Inches(2.3 + row * 1.4)
    box_w = Inches(3.1) if col == 0 else Inches(3.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(20)
    p2.font.color.rgb = DARK_BLUE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER

# ===== SLIDE 3: 晴天時チーム分け =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "チーム分け（晴天時）")

add_table(slide,
    ["チーム", "配線指揮者", "サブリーダー"],
    [
        ["A", "井上英明", "ー"],
        ["B", "坪内創", "小日向風磨"],
        ["C", "沓掛正太郎", "ー"],
        ["D", "前多綾美", "ー"],
        ["E", "斎藤快", "大久保風都"],
        ["F", "冨田雄平", "ー"],
        ["G", "山本力也", "ー"],
    ],
    L, Inches(1.2), CW)

# 必要物品 table below
add_text_box(slide, "必要物品（晴天時）", L, Inches(5.0), CW, Inches(0.5),
             font_size=18, bold=True, color=DARK_BLUE)
add_table(slide,
    ["チーム", "30m電ドラ", "50m電ドラ", "ブレーカータップ"],
    [
        ["A", "4", "0", "4"],
        ["B", "3", "2", "5"],
        ["C", "3", "0", "3"],
        ["D", "1", "2", "3"],
        ["E", "3", "2", "5"],
        ["F", "1", "0", "1"],
        ["G", "2", "2", "4"],
        ["合計", "17", "8", "25"],
    ],
    L, Inches(5.5), CW)

# ===== SLIDE 4: 雨天時チーム分け =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "チーム分け（雨天時）")

add_table(slide,
    ["チーム", "配線指揮者", "サブリーダー"],
    [
        ["A", "井上英明、坪内創", "小日向風磨"],
        ["B", "沓掛正太郎、前多綾美", "ー"],
        ["C", "斎藤快、冨田雄平", "大久保風都"],
        ["D", "山本力也", "ー"],
    ],
    L, Inches(1.2), CW)

add_text_box(slide, "必要物品（雨天時）", L, Inches(3.8), CW, Inches(0.5),
             font_size=18, bold=True, color=DARK_BLUE)
add_table(slide,
    ["チーム", "30m電ドラ", "50m電ドラ", "ブレーカータップ"],
    [
        ["A", "6", "2", "8"],
        ["B", "2", "4", "6"],
        ["C", "2", "0", "2"],
        ["D", "1", "1", "2"],
        ["合計", "11", "7", "18"],
    ],
    L, Inches(4.3), CW)

# ===== SLIDE 5: 必要物品リスト =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "必要物品")

items = [
    "ガムテープ（EGG）",
    "養生テープ（屋内の配線用・EGG）",
    "電工ドラム",
    "ゴミ袋",
    "ブレーカータップ",
]
add_bullet_list(slide, items, L, Inches(1.3), CW, Inches(4), font_size=20)

# ===== SLIDE 6: タイムスケジュール =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "タイムスケジュール")

schedule = [
    ("18:00〜18:10", "チーム分け及び\n配線場所への移動"),
    ("18:10〜18:30", "各グループで\n配線方法のレクチャー"),
    ("18:30〜22:00", "配線作業"),
]
for i, (time_str, desc) in enumerate(schedule):
    y = Inches(1.3 + i * 1.6)
    # Time box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, L, y, CW, Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = time_str
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Description below
    add_text_box(slide, desc, Inches(MARGIN + 0.3), y + Inches(0.7), Inches(CONTENT_W - 0.3), Inches(0.7),
                 font_size=16, color=BLACK)

# Flow image
add_text_box(slide, "配線当日の流れ", L, Inches(6.5), CW, Inches(0.5),
             font_size=16, bold=True, color=DARK_BLUE)
add_image_safe(slide, "image2.png", L, Inches(7.0), width=CW)

# ===== SLIDE 7: 配線箇所（晴天時）=====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "配線箇所（晴天時）")

add_accent_box(slide,
    "作業時には通信状態が悪くなる可能性があるため、\nあらかじめ資料をダウンロードしておくこと。\n(07_44th_配線テープ箇所（軽量版）.pdf)",
    L, Inches(1.2), CW, Inches(1.0), font_size=12, bg_color=WARN_BG)

add_table(slide,
    ["グループ", "配線番号", "参照資料"],
    [
        ["A", "講義棟エリア\n（①〜④）", "02_44th_講義棟エリア\n詳細配線図.pdf"],
        ["B", "物材棟エリア\n（⑤〜⑨）", "03_44th_物財・図書館エリア\n詳細配線図.pdf"],
        ["C", "図書館棟エリア\n（⑩〜⑫）", "04_44th_電気棟エリア\n詳細配線図.pdf"],
        ["D", "電気棟エリア\n（⑬〜⑮）", "04_44th_電気棟エリア\n詳細配線図.pdf"],
        ["E", "機械棟エリア\n（⑯〜⑱）", "05_44th_機械棟エリア\n詳細配線図.pdf"],
        ["F", "福利棟エリア\n（⑲）", "04_44th_電気棟エリア\n詳細配線図.pdf"],
        ["G", "体育館エリア\n（⑳〜㉒）", "06_44th_体育館エリア\n詳細配線図.pdf"],
    ],
    L, Inches(2.5), CW, row_height=Inches(0.55))

# ===== SLIDE 8: 配線箇所（雨天時）=====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "配線箇所（雨天時）")

add_accent_box(slide,
    "雨により屋外の配線が行えない場合、\n9月11日(木)は屋内部分のみの配線を行う。\n"
    "作業中に乾いた場合、追加で配線を行う。\nその判断は総指揮が行う。",
    L, Inches(1.2), CW, Inches(1.2), font_size=12, bg_color=WARN_BG)

add_table(slide,
    ["グループ", "配線番号", "参照資料"],
    [
        ["A", "講義棟エリア\n物材棟エリア", "02_44th_講義棟エリア詳細配線図.pdf\n03_44th_物財・図書館エリア詳細配線図.pdf"],
        ["B", "電気棟エリア", "04_44th_電気棟エリア詳細配線図.pdf"],
        ["C", "機械棟エリア", "05_44th_機械棟エリア詳細配線図.pdf"],
        ["D", "福利棟エリア", "04_44th_電気棟エリア詳細配線図.pdf\n06_44th_体育館エリア詳細配線図.pdf"],
    ],
    L, Inches(2.7), CW, row_height=Inches(0.55))

# Rain team maps - 2x2 grid
for i, (name, img) in enumerate([
    ("Aチーム", "image7.png"), ("Bチーム", "image25.png"),
    ("Cチーム", "image20.png"), ("Dチーム", "image1.png")
]):
    col = i % 2
    row = i // 2
    x = Inches(MARGIN + col * 3.5)
    y = Inches(5.5 + row * 3.5)
    add_text_box(slide, name, x, y, Inches(3.2), Inches(0.4),
                 font_size=13, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    add_image_safe(slide, img, x, y + Inches(0.4), width=Inches(3.2))

# ===== SLIDE 9: 配線時の注意点 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "配線時の注意点")

warnings = [
    "電源ドラムをコンセントに繋ぐ時は必ず配線指揮者が立ち合い、ブレーカータップを使用する。",
    "屋内部分の配線は養生テープにより配線を行う。",
    "配線作業について不明な点があれば必ず指揮担当者に確認を取る。",
    "雨により屋外の配線が行えないため屋内のみを配線する場合がある。",
    "作業中に配線箇所の地面が乾いた場合、その箇所の配線を行う。配線の指示は総指揮が行う。",
]
for i, w in enumerate(warnings):
    y = Inches(1.3 + i * 1.5)
    # Warning icon
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, L, y, Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "!"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, w, Inches(MARGIN + 0.7), y, Inches(CONTENT_W - 0.7), Inches(1.2), font_size=14, color=BLACK)

# ===== SLIDE 10: 配線計画図 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "配線計画図")

# 2x3 grid for vertical layout
config_images = ["image17.png", "image23.png", "image4.png", "image11.png", "image22.png", "image24.png"]
img_w = Inches(CONTENT_W / 2 - 0.1)
for i, img in enumerate(config_images):
    col = i % 2
    row = i // 2
    x = Inches(MARGIN + col * (CONTENT_W / 2 + 0.05))
    y = Inches(1.2 + row * 3.8)
    add_image_safe(slide, img, x, y, width=img_w)

# ===== SLIDE 11: 電源付近の処理 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "電源付近の処理")

add_text_box(slide,
    "コンセントにさして、プラグが取れないように\nその周りをガムテープで固定する。\n\n"
    "ガムテープは本部から持ってくる。\n\n"
    "固定後は、少し引っ張りプラグが\n取れないか確認する。",
    L, Inches(1.2), CW, Inches(2.5), font_size=16)

add_text_box(slide, "コンセント周りの固定", L, Inches(3.8), CW, Inches(0.5),
             font_size=15, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_image_safe(slide, "image9.jpg", L, Inches(4.3), width=CW)

# ===== SLIDE 12: コードの固定処理 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "コードの固定処理")

add_text_box(slide,
    "50cm程度の長さに切ったガムテープを使用して、\nコードの片側を固定する。\n"
    "同程度の長さのガムテープでもう片側を固定し、\n配線先まで同じ工程を繰り返す。\n\n"
    "※どの箇所もコードが裸で見えないように注意する。",
    L, Inches(1.2), CW, Inches(2.0), font_size=14)

# Step images stacked vertically
steps = [
    ("Step 1: ケーブル固定前", "image10.jpg"),
    ("Step 2: ケーブル固定途中", "image5.png"),
    ("Step 3: ケーブル固定後", "image19.png"),
    ("Step 4: 中央の固定", "image14.png"),
]
for i, (label, img) in enumerate(steps):
    y = Inches(3.3 + i * 2.5)
    add_text_box(slide, label, L, y, CW, Inches(0.4),
                 font_size=13, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_image_safe(slide, img, L, y + Inches(0.4), width=CW)

# ===== SLIDE 13: コードの固定場所 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "コードの固定場所")

add_text_box(slide,
    "通路の通行の邪魔にならないようにするため、\n配線は原則壁側（屋内の場合は幅木）に沿って行うこと。",
    L, Inches(1.2), CW, Inches(1.0), font_size=15)

fix_imgs = [
    ("壁側の固定", "image8.png"),
    ("床側の固定", "image15.jpg"),
    ("角部の固定", "image3.jpg"),
]
for i, (label, img) in enumerate(fix_imgs):
    y = Inches(2.5 + i * 3.5)
    add_text_box(slide, label, L, y, CW, Inches(0.4),
                 font_size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
    add_image_safe(slide, img, L, y + Inches(0.4), width=CW)

# ===== SLIDE 14: 自動ドア付近の配線 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "自動ドア付近・屋外の配線")

add_text_box(slide,
    "● 自動ドアではドアレール部分の中心にケーブルを通す。\n\n"
    "● 当日も自動ドアは稼働するため、ドアレール部分にはガムテープで固定しないこと。\n\n"
    "● 自動ドアにケーブルが3本以上通る場合は、ケーブルをピラミッド状にして配線を行う。\n\n"
    "● 屋外の配線も屋内同様に原則壁側に沿って行うこと。\n\n"
    "● タイルの溝がある箇所は壁側に近いくぼみを利用するなど工夫すること。\n\n"
    "● 自動ドア付近にマットがある場合は一度マットをよけてケーブルを固定し、固定後にマットを敷き直す。",
    L, Inches(1.2), CW, Inches(4.0), font_size=13)

add_text_box(slide, "自動ドア付近の配線", L, Inches(5.5), CW, Inches(0.4),
             font_size=14, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
# Two images side by side
half_w = Inches(CONTENT_W / 2 - 0.1)
add_image_safe(slide, "image16.jpg", L, Inches(6.0), width=half_w)
add_image_safe(slide, "image13.jpg", Inches(MARGIN + CONTENT_W / 2 + 0.05), Inches(6.0), width=half_w)
add_image_safe(slide, "image12.jpg", L, Inches(8.5), width=CW)

# ===== SLIDE 15: 模擬店テントスペースでの処理 =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "模擬店テントスペースでのケーブル処理")

add_text_box(slide,
    "● 天候に関係なく、供給先の模擬店テントスペースに到達したら、"
    "ドラムのケーブルが余っていた場合はすべて引き出すこと。\n\n"
    "● 雨天時や配線が余った時、テントのない箇所やまだテントが立てられていないスペースでは"
    "袋にドラムと余ったケーブルを入れる。",
    L, Inches(1.2), CW, Inches(2.0), font_size=14)

add_text_box(slide, "模擬店スペース内", L, Inches(3.5), CW, Inches(0.4),
             font_size=14, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_image_safe(slide, "image18.jpg", L, Inches(3.9), width=CW)

add_text_box(slide, "雨天時のドラム部の処理", L, Inches(7.5), CW, Inches(0.4),
             font_size=14, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
add_image_safe(slide, "image21.jpg", L, Inches(7.9), width=half_w)
add_image_safe(slide, "image6.jpg", Inches(MARGIN + CONTENT_W / 2 + 0.05), Inches(7.9), width=half_w)

# SLIDE 16（代表連絡先）はプライバシー保護のため除外

prs.save(OUT)
print(f"Saved: {OUT}")
