#!/usr/bin/env python3
"""汎用 docx → pptx 変換スクリプト（SKILL.md準拠）

Usage:
    python3 scripts/convert_docx.py docs/manuals/対象.docx [--work-dir scripts/manuals/短縮名]

- docxを読み取り、SKILL.mdのパターン選択基準に従って縦長pptxを自動生成
- 画像は --work-dir/images/ に抽出
- pptxは docs/manuals/ に出力
"""

import argparse
import os
import re
import sys
import zipfile

import docx
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── カラー定数（SKILL.md準拠）──────────────────────
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
TABLE_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)
LIGHT_BLUE_BG = RGBColor(0xD6, 0xE4, 0xF0)
WARN_BG = RGBColor(0xFF, 0xF3, 0xCD)

W = 7.5
MARGIN = 0.4
CONTENT_W = W - MARGIN * 2  # 6.7
L = Inches(MARGIN)
CW = Inches(CONTENT_W)

# ── 除外パターン ─────────────────────────────
PHONE_RE = re.compile(r'0\d{1,4}[-‐ー]?\d{1,4}[-‐ー]?\d{2,4}')
EMAIL_RE = re.compile(r'[\w.+-]+@[\w-]+\.[\w.]+')
TIME_RANGE_RE = re.compile(r'\d{1,2}:\d{2}\s*[〜~～ー−-]\s*\d{1,2}:\d{2}')
KV_RE = re.compile(r'^(.+?)[：:]\s*(.+)$')


def should_exclude_section(heading_text):
    """代表連絡先セクションなどを除外"""
    exclude_keywords = ['代表連絡先', '連絡先一覧', '緊急連絡先']
    return any(kw in heading_text for kw in exclude_keywords)


def strip_private(text):
    """電話番号・メールアドレスを除去"""
    text = PHONE_RE.sub('（非公開）', text)
    text = EMAIL_RE.sub('（非公開）', text)
    return text


# ── docx 読み取り ─────────────────────────────

def extract_images(docx_path, images_dir):
    """docxから画像を抽出"""
    os.makedirs(images_dir, exist_ok=True)
    with zipfile.ZipFile(docx_path, 'r') as z:
        for name in z.namelist():
            if name.startswith('word/media/'):
                fname = os.path.basename(name)
                with open(os.path.join(images_dir, fname), 'wb') as f:
                    f.write(z.read(name))


NSMAP = {
    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
}


def get_para_text(elem):
    """パラグラフ要素からテキストを抽出（SDT対応）"""
    texts = elem.findall('.//w:t', NSMAP)
    return ''.join(t.text or '' for t in texts).strip()


def get_para_style(elem):
    """パラグラフ要素からスタイル名を取得"""
    pPr = elem.find('.//w:pStyle', NSMAP)
    if pPr is not None:
        val = pPr.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', '')
        # 'Heading1' → 'Heading 1', 'Title' → 'Title'
        import re as _re
        m = _re.match(r'(Heading)(\d+)', val)
        if m:
            return f'{m.group(1)} {m.group(2)}'
        return val
    return 'normal'


def get_para_images(elem, doc):
    """パラグラフ要素から画像ファイル名を抽出"""
    images = []
    drawings = elem.findall('.//w:drawing', NSMAP)
    for d in drawings:
        blips = d.findall('.//a:blip', NSMAP)
        for blip in blips:
            rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if rId and rId in doc.part.rels:
                target = doc.part.rels[rId].target_ref
                images.append(os.path.basename(target))
    return images


def build_image_map(doc):
    """パラグラフindex → 画像ファイル名 のマッピングを構築"""
    img_map = {}
    body = doc.element.body
    para_idx = 0
    for child in body:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'p':
            imgs = get_para_images(child, doc)
            if imgs:
                img_map[para_idx] = imgs
            para_idx += 1
    return img_map


def parse_docx_sections(doc, img_map):
    """docxをセクション（Heading1区切り）に分割して構造化
    XML要素を直接読み取ることでSDT内テキストにも対応"""
    sections = []
    current_section = None
    current_subsection = None

    body = doc.element.body
    para_idx = 0
    tbl_idx = 0

    title_text = None
    subtitle_text = None

    for child in body:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag

        if tag == 'tbl':
            # テーブル: doc.tables[tbl_idx] からデータ取得
            if tbl_idx < len(doc.tables):
                tbl = doc.tables[tbl_idx]
                # SDT対応: セルのテキストもXMLから直接取得
                headers = []
                rows_data = []
                for r_idx, row in enumerate(tbl.rows):
                    cells_text = []
                    for cell in row.cells:
                        # cell.text が空の場合、XMLから直接取得
                        ct = cell.text.strip()
                        if not ct:
                            t_elems = cell._element.findall('.//w:t', NSMAP)
                            ct = ''.join(t.text or '' for t in t_elems).strip()
                        cells_text.append(strip_private(ct))
                    if r_idx == 0:
                        headers = cells_text
                    else:
                        rows_data.append(cells_text)

                table_data = {'headers': headers, 'rows': rows_data}
                if current_subsection is not None:
                    current_subsection.setdefault('tables', []).append(table_data)
                elif current_section is not None:
                    current_section.setdefault('tables', []).append(table_data)
            tbl_idx += 1
            continue

        if tag == 'sdt':
            # SDT要素内のパラグラフを処理
            sdt_paras = child.findall('.//w:p', NSMAP)
            for sp in sdt_paras:
                # SDT内のパラグラフもカウント（doc.paragraphsとの整合性）
                pass
            # SDT全体からテキストを取得して現在のセクションに追加
            sdt_text = get_para_text(child).strip()
            if sdt_text and current_section and not current_section.get('excluded'):
                sdt_text = strip_private(sdt_text)
                target = current_subsection if current_subsection else current_section
                if target:
                    target.setdefault('paragraphs', []).append(sdt_text)
            continue

        if tag != 'p':
            continue

        # パラグラフ処理（XML直接読み取り）
        style = get_para_style(child)
        text = get_para_text(child)
        images = get_para_images(child, doc)

        if not text and not images:
            para_idx += 1
            continue

        # Title → タイトルスライド
        if style == 'Title':
            if title_text is None:
                title_text = strip_private(text)
            else:
                subtitle_text = strip_private(text)
            para_idx += 1
            continue

        # Heading 1 → 新セクション
        if style.startswith('Heading 1'):
            if should_exclude_section(text):
                current_section = {'heading': text, 'excluded': True, 'subsections': [], 'paragraphs': [], 'images': []}
            else:
                current_section = {'heading': strip_private(text), 'subsections': [], 'paragraphs': [], 'images': [], 'tables': []}
            sections.append(current_section)
            current_subsection = None
            para_idx += 1
            continue

        # Heading 2+ → サブセクション
        if style.startswith('Heading'):
            if current_section and not current_section.get('excluded'):
                current_subsection = {'heading': strip_private(text), 'paragraphs': [], 'images': [], 'tables': []}
                current_section['subsections'].append(current_subsection)
            para_idx += 1
            continue

        # 除外セクション内はスキップ
        if current_section and current_section.get('excluded'):
            para_idx += 1
            continue

        # 画像
        if images:
            target = current_subsection if current_subsection else current_section
            if target:
                target.setdefault('images', []).extend(images)

        # テキスト
        if text:
            text = strip_private(text)
            target = current_subsection if current_subsection else current_section
            if target:
                target.setdefault('paragraphs', []).append(text)

        para_idx += 1

    return title_text, subtitle_text, sections


# ── パターン分類 ──────────────────────────────

def classify_content(subsection):
    """サブセクションの内容からレイアウトパターンを判定"""
    paragraphs = subsection.get('paragraphs', [])
    tables = subsection.get('tables', [])
    images = subsection.get('images', [])
    heading = subsection.get('heading', '')

    # テーブルがある → C
    if tables:
        return 'C'

    # 画像がある → F（写真ステップ）
    if images:
        return 'F'

    # key-value形式が2つ以上 → B（情報カード）
    kv_count = sum(1 for p in paragraphs if KV_RE.match(p))
    if 2 <= kv_count <= 6:
        return 'B'

    # タイムライン（HH:MM〜HH:MM）
    timeline_count = sum(1 for p in paragraphs if TIME_RANGE_RE.search(p))
    if timeline_count >= 2:
        return 'D'

    # アラート（「必ず」「〜すること」「禁止」）
    alert_keywords = ['必ず', 'すること', '禁止', '注意', '厳禁', '絶対']
    alert_count = sum(1 for p in paragraphs if any(kw in p for kw in alert_keywords))
    if alert_count >= 2:
        return 'E'

    # それ以外 → G（箇条書き）
    return 'G'


# ── pptx ヘルパー関数 ──────────────────────────

def add_bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(prs, slide, text, y=Inches(0.2), height=Inches(0.7)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), y, prs.slide_width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_text_box(slide, text, left, top, width, height, font_size=16, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=18):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "● " + item if not item.startswith("●") else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8)
    return txBox


def add_image_safe(slide, img_path, left, top, width=None, height=None):
    if os.path.exists(img_path):
        kwargs = {"image_file": img_path, "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        try:
            return slide.shapes.add_picture(**kwargs)
        except Exception as e:
            print(f"  Warning: Could not add image {img_path}: {e}")
    return None


def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.4)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    if n_cols == 0:
        return None
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_height * n_rows)
    table = table_shape.table
    col_w = int(width / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_w
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            if c >= n_cols:
                break
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = BLACK
                p.alignment = PP_ALIGN.CENTER
    return table_shape


def add_info_card(slide, label, value, x, y, w=Inches(3.1), h=Inches(1.1)):
    """パターンB: 情報カード"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_BLUE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(20)
    p2.font.color.rgb = DARK_BLUE
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER


def add_timeline_item(slide, time_str, desc, y):
    """パターンD: タイムラインアイテム"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, L, y, CW, Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = time_str
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_text_box(slide, desc, Inches(MARGIN + 0.3), y + Inches(0.7),
                 Inches(CONTENT_W - 0.3), Inches(0.7), font_size=16, color=BLACK)


def add_alert_item(slide, text, y):
    """パターンE: アラートアイテム"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, L, y, Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ORANGE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "!"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    add_text_box(slide, text, Inches(MARGIN + 0.7), y,
                 Inches(CONTENT_W - 0.7), Inches(1.2), font_size=14, color=BLACK)


# ── スライド生成 ──────────────────────────────

def create_title_slide(prs, title, subtitle=None):
    """パターンA: タイトルスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BLUE)
    add_text_box(slide, title, L, Inches(4), CW, Inches(1.5),
                 font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, subtitle, L, Inches(5.5), CW, Inches(0.8),
                     font_size=24, color=RGBColor(0xA0, 0xC4, 0xE8), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, "第44回 技大祭", L, Inches(6.3), CW, Inches(0.8),
                 font_size=18, color=RGBColor(0xA0, 0xC4, 0xE8), alignment=PP_ALIGN.CENTER)
    return slide


def create_info_card_slide(prs, title, kv_pairs):
    """パターンB: 情報カードスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(prs, slide, title)
    for i, (label, value) in enumerate(kv_pairs[:6]):
        col = i % 2
        row = i // 2
        x = Inches(MARGIN + col * 3.5)
        y = Inches(1.3 + row * 1.4)
        add_info_card(slide, label, value, x, y)
    return slide


def create_table_slides(prs, title, tables, images_dir):
    """パターンC: テーブルスライド（8行で分割）"""
    slides = []
    for tbl_data in tables:
        headers = tbl_data['headers']
        all_rows = tbl_data['rows']
        # 8行ずつ分割
        chunk_size = 8
        for chunk_start in range(0, max(len(all_rows), 1), chunk_size):
            chunk = all_rows[chunk_start:chunk_start + chunk_size]
            if not chunk:
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            slide_title = title
            if len(all_rows) > chunk_size:
                page = chunk_start // chunk_size + 1
                total = (len(all_rows) + chunk_size - 1) // chunk_size
                slide_title = f"{title}（{page}/{total}）"
            add_title_bar(prs, slide, slide_title)
            add_table(slide, headers, chunk, L, Inches(1.2), CW)
            slides.append(slide)
    return slides


def create_timeline_slide(prs, title, items):
    """パターンD: タイムラインスライド"""
    # 最大6項目/スライド
    slides = []
    chunk_size = 6
    for chunk_start in range(0, len(items), chunk_size):
        chunk = items[chunk_start:chunk_start + chunk_size]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_title_bar(prs, slide, title)
        for i, (time_str, desc) in enumerate(chunk):
            y = Inches(1.3 + i * 1.6)
            add_timeline_item(slide, time_str, desc, y)
        slides.append(slide)
    return slides


def create_alert_slide(prs, title, warnings):
    """パターンE: アラートスライド"""
    slides = []
    chunk_size = 5
    for chunk_start in range(0, len(warnings), chunk_size):
        chunk = warnings[chunk_start:chunk_start + chunk_size]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_title_bar(prs, slide, title)
        for i, w in enumerate(chunk):
            y = Inches(1.3 + i * 1.5)
            add_alert_item(slide, w, y)
        slides.append(slide)
    return slides


def create_photo_slide(prs, title, images, images_dir):
    """パターンF: 写真ステップスライド（最大2枚/スライド）"""
    slides = []
    chunk_size = 2
    for chunk_start in range(0, len(images), chunk_size):
        chunk = images[chunk_start:chunk_start + chunk_size]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_title_bar(prs, slide, title)
        for i, img_name in enumerate(chunk):
            y = Inches(1.3 + i * 5.5)
            step_label = f"Step {chunk_start + i + 1}"
            add_text_box(slide, step_label, L, y, CW, Inches(0.4),
                         font_size=13, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
            img_path = os.path.join(images_dir, img_name)
            add_image_safe(slide, img_path, L, y + Inches(0.4), width=CW)
        slides.append(slide)
    return slides


def create_bullet_slide(prs, title, items):
    """パターンG: 箇条書きスライド"""
    slides = []
    # 200文字以内を目安に分割
    chunk = []
    char_count = 0
    for item in items:
        if char_count + len(item) > 200 and chunk:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_title_bar(prs, slide, title)
            add_bullet_list(slide, chunk, L, Inches(1.3), CW, Inches(10))
            slides.append(slide)
            chunk = []
            char_count = 0
        chunk.append(item)
        char_count += len(item)
    if chunk:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_title_bar(prs, slide, title)
        add_bullet_list(slide, chunk, L, Inches(1.3), CW, Inches(10))
        slides.append(slide)
    return slides


def create_mixed_slide(prs, title, paragraphs, tables, images, images_dir):
    """テキスト+テーブル+画像が混在するセクション用"""
    slides = []
    # テーブルがあればテーブルスライド
    if tables:
        slides.extend(create_table_slides(prs, title, tables, images_dir))
    # 画像があれば写真スライド
    if images:
        slides.extend(create_photo_slide(prs, title, images, images_dir))
    # テキストのみ（テーブル・画像がない場合）
    if paragraphs and not tables and not images:
        slides.extend(create_bullet_slide(prs, title, paragraphs))
    return slides


# ── メイン処理 ────────────────────────────────

def process_section(prs, section, images_dir):
    """セクション（Heading1配下）をスライド化"""
    if section.get('excluded'):
        return

    heading = section.get('heading', '')
    subsections = section.get('subsections', [])
    paragraphs = section.get('paragraphs', [])
    tables = section.get('tables', [])
    images = section.get('images', [])

    if subsections:
        # サブセクションごとに処理
        # まずセクション直下のコンテンツ
        if paragraphs or tables or images:
            # セクション直下の key-value → 情報カード
            kv_pairs = []
            other_paras = []
            for p in paragraphs:
                m = KV_RE.match(p)
                if m:
                    kv_pairs.append((m.group(1).strip(), m.group(2).strip()))
                else:
                    other_paras.append(p)
            if 2 <= len(kv_pairs) <= 6:
                create_info_card_slide(prs, heading, kv_pairs)
            elif other_paras:
                create_bullet_slide(prs, heading, other_paras)
            if tables:
                create_table_slides(prs, heading, tables, images_dir)

        for sub in subsections:
            process_subsection(prs, sub, heading, images_dir)
    else:
        # サブセクションなし → セクション全体を1つのコンテンツとして処理
        all_content = {'heading': heading, 'paragraphs': paragraphs, 'tables': tables, 'images': images}
        pattern = classify_content(all_content)
        generate_slide_by_pattern(prs, pattern, heading, all_content, images_dir)


def process_subsection(prs, sub, parent_heading, images_dir):
    """サブセクション（Heading2配下）をスライド化"""
    sub_heading = sub.get('heading', '')
    title = sub_heading or parent_heading

    pattern = classify_content(sub)
    generate_slide_by_pattern(prs, pattern, title, sub, images_dir)


def generate_slide_by_pattern(prs, pattern, title, content, images_dir):
    """パターンに応じたスライドを生成"""
    paragraphs = content.get('paragraphs', [])
    tables = content.get('tables', [])
    images = content.get('images', [])

    if not paragraphs and not tables and not images:
        return

    if pattern == 'B':
        kv_pairs = []
        for p in paragraphs:
            m = KV_RE.match(p)
            if m:
                kv_pairs.append((m.group(1).strip(), m.group(2).strip()))
        if kv_pairs:
            create_info_card_slide(prs, title, kv_pairs)
        # テーブルがあれば追加
        if tables:
            create_table_slides(prs, title, tables, images_dir)

    elif pattern == 'C':
        # テーブル前のテキストがあれば箇条書きで出す
        non_table_text = [p for p in paragraphs if not p.startswith('表')]
        if non_table_text:
            create_bullet_slide(prs, title, non_table_text)
        create_table_slides(prs, title, tables, images_dir)
        if images:
            create_photo_slide(prs, title, images, images_dir)

    elif pattern == 'D':
        items = []
        for p in paragraphs:
            m = TIME_RANGE_RE.search(p)
            if m:
                time_str = m.group(0)
                # 時間の後の説明を抽出
                rest = p[m.end():].strip().lstrip('：:： ').strip()
                if not rest:
                    # "HH:MM〜HH:MM：説明" 形式
                    parts = p.split('：', 1) if '：' in p else p.split(':', 1)
                    if len(parts) > 1:
                        rest = parts[-1].strip()
                items.append((time_str, rest))
            else:
                # タイムライン以外のテキストは前の項目の説明に追加
                if items:
                    prev_time, prev_desc = items[-1]
                    items[-1] = (prev_time, prev_desc + '\n' + p if prev_desc else p)
        if items:
            create_timeline_slide(prs, title, items)
        if tables:
            create_table_slides(prs, title, tables, images_dir)

    elif pattern == 'E':
        create_alert_slide(prs, title, paragraphs)
        if tables:
            create_table_slides(prs, title, tables, images_dir)

    elif pattern == 'F':
        # テキストがあれば先に表示
        if paragraphs:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_title_bar(prs, slide, title)
            text = '\n'.join(paragraphs[:5])
            add_text_box(slide, text, L, Inches(1.2), CW, Inches(2.0), font_size=14)
            # 最初の2枚の画像をこのスライドに配置
            for i, img_name in enumerate(images[:2]):
                y = Inches(3.5 + i * 4.5)
                img_path = os.path.join(images_dir, img_name)
                add_image_safe(slide, img_path, L, y, width=CW)
            # 残りの画像は別スライド
            if len(images) > 2:
                create_photo_slide(prs, title, images[2:], images_dir)
        else:
            create_photo_slide(prs, title, images, images_dir)
        if tables:
            create_table_slides(prs, title, tables, images_dir)

    elif pattern == 'G':
        create_bullet_slide(prs, title, paragraphs)
        if tables:
            create_table_slides(prs, title, tables, images_dir)
        if images:
            create_photo_slide(prs, title, images, images_dir)


def convert_docx_to_pptx(docx_path, work_dir, output_dir):
    """メイン: docx → pptx 変換"""
    images_dir = os.path.join(work_dir, 'images')
    docx_name = os.path.splitext(os.path.basename(docx_path))[0]
    pptx_path = os.path.join(output_dir, docx_name + '.pptx')

    print(f"Converting: {os.path.basename(docx_path)}")
    print(f"  Work dir: {work_dir}")
    print(f"  Output:   {pptx_path}")

    # 1. 画像抽出
    extract_images(docx_path, images_dir)
    img_count = len([f for f in os.listdir(images_dir) if not f.startswith('.')])
    print(f"  Extracted {img_count} images")

    # 2. docx読み取り
    doc = docx.Document(docx_path)
    img_map = build_image_map(doc)
    title, subtitle, sections = parse_docx_sections(doc, img_map)

    print(f"  Title: {title}")
    print(f"  Sections: {len(sections)}")

    # 3. pptx生成
    prs = Presentation()
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(13.333)

    # タイトルスライド
    if title:
        create_title_slide(prs, title, subtitle)

    # 各セクション
    for section in sections:
        process_section(prs, section, images_dir)

    # 保存
    os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
    prs.save(pptx_path)
    slide_count = len(prs.slides)
    print(f"  Generated {slide_count} slides → {pptx_path}")
    return pptx_path


def derive_short_name(docx_filename):
    """docxファイル名から短縮名を導出"""
    name = os.path.splitext(docx_filename)[0]
    # プレフィックス除去: 01_44th_, 44th_, 02_44th_, 44th_06_ など
    name = re.sub(r'^\d+_?(44th_)?(\d+_)?', '', name)
    # 「マニュアル」以降を除去
    name = re.sub(r'マニュアル.*$', '', name)
    # 残りをクリーンアップ
    name = name.strip('_')
    if not name:
        name = os.path.splitext(docx_filename)[0][:20]
    return name


def main():
    parser = argparse.ArgumentParser(description='docx → pptx 変換（SKILL.md準拠）')
    parser.add_argument('docx_path', help='入力docxファイルパス')
    parser.add_argument('--work-dir', help='作業ディレクトリ（省略時は自動生成）')
    parser.add_argument('--output-dir', default=None, help='pptx出力先（省略時はdocs/manuals/）')
    args = parser.parse_args()

    docx_path = os.path.abspath(args.docx_path)
    if not os.path.exists(docx_path):
        print(f"Error: {docx_path} not found")
        sys.exit(1)

    # プロジェクトルート
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_root = os.path.dirname(script_dir)

    # 作業ディレクトリ
    if args.work_dir:
        work_dir = os.path.abspath(args.work_dir)
    else:
        short_name = derive_short_name(os.path.basename(docx_path))
        work_dir = os.path.join(project_root, 'scripts', 'manuals', short_name)

    # 出力ディレクトリ
    output_dir = args.output_dir or os.path.join(project_root, 'docs', 'manuals')

    os.makedirs(work_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    convert_docx_to_pptx(docx_path, work_dir, output_dir)


if __name__ == '__main__':
    main()
